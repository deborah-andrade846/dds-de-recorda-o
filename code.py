# -*- coding: utf-8 -*-
"""
Editor Spyder

Este é um arquivo de script temporário.
"""

from pptx import Presentation
from pptx.util import Pt

# 1. Criar a apresentação
prs = Presentation()

# 2. Criar o Slide de Título (Slide 1)
slide_title_layout = prs.slide_layouts[0] # 0 é o layout padrão de título
slide_title = prs.slides.add_slide(slide_title_layout)

title = slide_title.shapes.title
subtitle = slide_title.placeholders[1]

title.text = "DDS: Recordação de Acidente"
subtitle.text = "Atenção com Animais nas Portarias e Áreas Operacionais\nGincana de Segurança"

# Função auxiliar para criar slides de conteúdo com formatação
def add_content_slide(prs, title_text, bullet_points):
    slide_layout = prs.slide_layouts[1] # 1 é o layout de título e conteúdo
    slide = prs.slides.add_slide(slide_layout)
    
    # Configurar título
    title = slide.shapes.title
    title.text = title_text
    
    # Configurar corpo do texto
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear() # Limpa o texto padrão
    
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(24)
        p.level = 0 # Nível do bullet point

    return slide

# 3. Criar o Slide: O que aconteceu? (Slide 2)
pontos_fato = [
    "Data: 10/06/2026 às 11:50 | Local: Portaria.",
    "Fato: Colaborador retornava ao veículo após passar pela catraca e foi surpreendido por um cachorro.",
    "O colaborador sofreu uma mordida na panturrilha esquerda.",
    "Classificação: Quase Acidente (Nível 1 - Leve).",
    "Consequência: Apenas arranhão, sem necessidade de intervenção médica."
]
add_content_slide(prs, "O que aconteceu? (Recapitulando)", pontos_fato)

# 4. Criar o Slide: Lições Aprendidas (Slide 3)
pontos_prevencao = [
    "Evite aproximação e contato: Mesmo que o animal pareça dócil, não tente passar a mão.",
    "Não alimente os animais: Isso faz com que eles permaneçam na área operacional.",
    "Atenção aos arredores: Ao caminhar, evite distrações como o uso do celular.",
    "Comunique imediatamente: Avise o Meio Ambiente, Segurança ou a Portaria sobre animais em local de risco."
]
add_content_slide(prs, "Lições Aprendidas & Como Prevenir", pontos_prevencao)

# 5. Criar o Slide: Gincana (Slide 4)
pontos_gincana = [
    "Lembrete essencial para a pontuação da Gincana de Segurança!",
    "Precisamos comprovar 100% de adesão nos turnos.",
    "Assine a lista de presença ESPECÍFICA da letra do seu turno (ex: uma lista só para Letra A, outra para Letra B).",
    "Não misture as assinaturas em uma lista única. Isso reduz a nossa nota.",
    "Vamos garantir os 10 pontos da nossa equipe!"
]
add_content_slide(prs, "Atenção: Assinaturas da Gincana", pontos_gincana)

# 6. Salvar o arquivo
nome_arquivo = 'DDS_Recordacao_Acidente.pptx'
prs.save(nome_arquivo)

print(f"Sucesso! A apresentação '{nome_arquivo}' foi gerada e está pronta para uso.")